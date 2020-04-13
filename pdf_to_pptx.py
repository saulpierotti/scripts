#!/home/saul/miniconda3/envs/pdf_to_pptx/bin/python

# Adapted by Saul from:
#    git@github.com:TomAnthony/pdf-to-powerpoint.git

# Requirements:
#   python3
#   lxml==4.2.1
#   pdf2image==0.1.14
#   Pillow==5.1.0
#   python-pptx==0.6.10
#   XlsxWriter==1.0.5
#   poppler


import os, sys

from PIL import Image
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

pdf_file = sys.argv[1]
print()
print("Converting file: " + pdf_file)
print()

# Prep presentation
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]

# Create working folder
base_name = pdf_file.split(".pdf")[0]

# Convert PDF to list of images
print("Starting conversion...")
slideimgs = convert_from_path(pdf_file, 300, fmt="ppm", thread_count=2)
print("...complete.")
print()

# Loop over slides
for i, slideimg in enumerate(slideimgs):
    if i % 10 == 0:
        print("Saving slide: " + str(i))

    imagefile = BytesIO()
    slideimg.save(imagefile, format="tiff")
    imagedata = imagefile.getvalue()
    imagefile.seek(0)
    width, height = slideimg.size

    # Set slide dimensions
    prs.slide_height = height * 9525
    prs.slide_width = width * 9525

    # Add slide
    slide = prs.slides.add_slide(blank_slide_layout)
    pic = slide.shapes.add_picture(
        imagefile, 0, 0, width=width * 9525, height=height * 9525
    )

# Save Powerpoint
print()
print("Saving file: " + base_name + ".pptx")
prs.save(base_name + ".pptx")
print("Conversion complete. :)")
print()
